from pptx import Presentation
from pptx.util import Pt

prs = Presentation()

def add_slide(title, points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()

    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(20)

# Slides
add_slide("Blogging System", [
    "Using Django, HTML, CSS & Bootstrap",
    "Full Stack Web Application",
    "Your Name"
])

add_slide("Introduction", [
    "Web application for blogging",
    "Users can create and manage posts",
    "Dynamic and interactive system"
])

add_slide("Objectives", [
    "Develop blogging platform",
    "Implement authentication",
    "Enable CRUD operations",
    "Responsive design"
])

add_slide("Technologies Used", [
    "Python (Django)",
    "HTML, CSS",
    "Bootstrap",
    "SQLite Database"
])

add_slide("Why Django?", [
    "Fast and secure framework",
    "Built-in admin panel",
    "ORM for database",
    "Scalable architecture"
])

add_slide("System Architecture", [
    "User interacts via browser",
    "Frontend sends request",
    "Django processes logic",
    "Database interaction",
    "Response displayed"
])

add_slide("Features", [
    "User login/register",
    "Create, edit, delete posts",
    "Responsive UI",
    "Admin dashboard"
])

add_slide("Project Structure", [
    "models.py",
    "views.py",
    "templates/",
    "static/",
    "urls.py"
])

add_slide("Database Design", [
    "Post: Title, Content, Author",
    "Created Date",
    "User model",
    "Comments (optional)"
])

add_slide("Frontend Design", [
    "HTML structure",
    "CSS styling",
    "Bootstrap for responsiveness"
])

add_slide("Functionalities", [
    "Authentication system",
    "Blog management",
    "Search feature",
    "Pagination"
])

add_slide("Advantages", [
    "Easy to use",
    "Secure",
    "Scalable",
    "Responsive"
])

add_slide("Limitations", [
    "Basic UI",
    "Limited features",
    "Requires deployment"
])

add_slide("Future Enhancements", [
    "Like and share feature",
    "REST API integration",
    "Tags and categories",
    "Cloud deployment"
])

add_slide("Conclusion", [
    "Built using Django",
    "Full-stack development",
    "Can be extended further"
])

add_slide("Thank You", [
    "Questions?"
])

# Save file
prs.save("Blogging_System_Presentation.pptx")

print("PPT created successfully!")